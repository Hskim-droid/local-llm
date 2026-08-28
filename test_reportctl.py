#!/usr/bin/env python3
"""같은 입력 → 같은 내용. 모델 없이 결정적 구간만 검증."""
from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path

from extract import Segment, collect_numbers, extract, extract_many, hangul_ratio
from hardware import pick_profile
from ollama_client import load_config
from reportctl import parse_drop_line
from render import build_html, render_md, write_outputs
from structure import default_title, structure


ROOT = Path(__file__).resolve().parent


class ExtractTests(unittest.TestCase):
    def test_pptx_slides_and_table(self):
        from pptx import Presentation
        from pptx.util import Inches, Pt

        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "deck.pptx"
            prs = Presentation()
            blank = prs.slide_layouts[6]
            s1 = prs.slides.add_slide(blank)
            box = s1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
            box.text_frame.paragraphs[0].text = "PartnerForum 2026 개요"
            box.text_frame.paragraphs[0].font.size = Pt(20)
            s2 = prs.slides.add_slide(blank)
            rows, cols = 2, 2
            table = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(1), Inches(8), Inches(1.5)).table
            table.cell(0, 0).text = "구분"
            table.cell(0, 1).text = "내용"
            table.cell(1, 0).text = "투자"
            table.cell(1, 1).text = "7억원"
            prs.save(path)
            segs = extract(path)
        locs = [s.location for s in segs]
        blob = "\n".join(s.text for s in segs)
        self.assertIn("slide-1", locs)
        self.assertIn("slide-2", locs)
        self.assertIn("PartnerForum 2026", blob)
        self.assertIn("7억원", blob)

    def test_txt_split(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "m.txt"
            path.write_text("1. 개요\n회사는 설립 2년됨.\n\n2. 논의\n후속 미팅 예정.\n", encoding="utf-8")
            segs = extract(path)
        self.assertGreaterEqual(len(segs), 1)
        self.assertTrue(any("설립 2년" in s.text for s in segs))

    def test_old_ppt_rejected(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "old.ppt"
            path.write_bytes(b"x")
            with self.assertRaises(ValueError):
                extract(path)

    def test_video_uses_sidecar_without_whisper(self):
        with tempfile.TemporaryDirectory() as tmp:
            vid = Path(tmp) / "회의.mp4"
            vid.write_bytes(b"not-a-real-video")
            vid.with_suffix(".txt").write_text(
                "[00:12] 시드 7억원을 논의함.\n[01:05] 9월 베타 배포 예정.\n",
                encoding="utf-8",
            )
            segs = extract(vid)
        blob = "\n".join(s.text for s in segs)
        self.assertTrue(segs)
        self.assertEqual(segs[0].source, "video")
        self.assertIn("7억원", blob)
        self.assertIn("9월", blob)

    def test_three_inputs_mixed(self):
        from pptx import Presentation
        from pptx.util import Inches
        import fitz

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            ppt = root / "장표.pptx"
            prs = Presentation()
            s = prs.slides.add_slide(prs.slide_layouts[6])
            box = s.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(8), Inches(1))
            box.text_frame.paragraphs[0].text = "장표 핵심 50억원"
            prs.save(ppt)
            pdf = root / "자료.pdf"
            doc = fitz.open()
            page = doc.new_page()
            page.insert_text((72, 72), "PDF 본문 hit rate 65%")
            doc.save(pdf)
            doc.close()
            vid = root / "회의.mp4"
            vid.write_bytes(b"x")
            vid.with_suffix(".txt").write_text("[00:01] 영상에서 후속 미팅 예정.\n", encoding="utf-8")
            segs = extract_many([vid, ppt, pdf])
        blob = "\n".join(s.text for s in segs)
        self.assertIn("50억원", blob)
        self.assertIn("65%", blob)
        self.assertIn("후속 미팅", blob)
        kinds = {s.source for s in segs}
        self.assertIn("video", kinds)
        self.assertIn("ppt", kinds)
        self.assertIn("pdf", kinds)


class NumberAndLangTests(unittest.TestCase):
    def test_numbers_kept(self):
        nums = collect_numbers("시드 7억원, hit rate 65%, 인원 3명")
        self.assertIn("7억원", nums)
        self.assertIn("65%", nums)
        self.assertIn("3명", nums)

    def test_hangul_ratio(self):
        self.assertGreater(hangul_ratio("회사는 설립 2년됨"), 0.7)
        self.assertLess(hangul_ratio("Company overview 2026"), 0.2)


class StructureRenderTests(unittest.TestCase):
    def setUp(self):
        self.segs = [
            Segment("예시바이오는 ApexBind를 소개함. 시드 7억원.", "txt", "part-1"),
            Segment("베타 버전은 9월 배포 예정.", "txt", "part-2"),
        ]

    def test_no_llm_schema_shape(self):
        d = structure(self.segs, title="예시바이오", kind="minutes", client=None, chunk_chars=1600)
        self.assertEqual(d["closing"], "끝.")
        self.assertEqual(len(d["header"]), 4)
        headings = [s["heading"] for s in d["sections"]]
        self.assertTrue(headings[0].startswith("0. Executive Summary"))
        self.assertTrue(any("향후 계획" in h for h in headings))
        dump = json.dumps(d, ensure_ascii=False)
        self.assertIn("7억원", dump)
        self.assertIn("9월", dump)

    def test_md_and_html_contain_form(self):
        d = structure(self.segs, title="예시바이오", kind="minutes", client=None, chunk_chars=1600)
        md = render_md(d)
        html = build_html(d)
        self.assertIn("Executive Summary", md)
        self.assertIn("끝.", md)
        self.assertIn("&lt; 회 의 내 용 &gt;", html)
        self.assertIn("예시바이오", html)

    def test_write_docx_is_final(self):
        d = structure(self.segs, title="테스트", kind="slides", client=None, chunk_chars=1600)
        with tempfile.TemporaryDirectory() as tmp:
            paths = write_outputs(d, Path(tmp), "보고서")
            self.assertIn("docx", paths)
            self.assertTrue(paths["docx"].is_file())
            self.assertNotIn("pdf", paths)
            self.assertNotIn("html", paths)
            from docx import Document
            doc = Document(str(paths["docx"]))
            text = "\n".join(p.text for p in doc.paragraphs)
            self.assertIn("테스트", text)
            self.assertIn("7억원", text)
            self.assertIn("끝.", text)
            self.assertNotIn("< 회 의 내 용 >", text)

    def test_default_title(self):
        self.assertEqual(default_title([Path("/tmp/PartnerForum_GPDay.pptx")]), "PartnerForum GPDay")


class ProfileTests(unittest.TestCase):
    def setUp(self):
        self.cfg = load_config(ROOT)

    def test_gram16(self):
        pid, p = pick_profile(self.cfg, ram=16, platform="win32")
        self.assertEqual(pid, "gram16")
        self.assertEqual(p["pull"], "qwen3:8b")
        self.assertEqual(p["setup_models"], ["qwen3:8b"])
        self.assertLessEqual(int(p["num_ctx"]), 4096)

    def test_gram32(self):
        pid, p = pick_profile(self.cfg, ram=32, platform="win32")
        self.assertEqual(pid, "gram32")
        self.assertIn("qwen3:14b", p["model_pref"])

    def test_force(self):
        pid, p = pick_profile(self.cfg, ram=32, explicit="gram16")
        self.assertEqual(pid, "gram16")


class DropLineTests(unittest.TestCase):
    def test_escaped_spaces(self):
        line = r"/tmp/회의\ 영상.mp4 /tmp/장표.pptx /tmp/자료.pdf"
        paths = parse_drop_line(line)
        self.assertEqual(len(paths), 3)
        self.assertEqual(paths[0].name, "회의 영상.mp4")
        self.assertEqual(paths[1].suffix, ".pptx")

    def test_quoted(self):
        paths = parse_drop_line('"/tmp/a b.mp4" "/tmp/c.pptx"')
        self.assertEqual(paths[0].name, "a b.mp4")

    def test_windows_quoted(self):
        paths = parse_drop_line(r'"C:\Users\kim\회의 영상.mp4" "C:\Users\kim\장표.pptx"')
        self.assertEqual(len(paths), 2)
        self.assertTrue(str(paths[0]).endswith("회의 영상.mp4"))


class ExampleGoldenTests(unittest.TestCase):
    def test_example_no_llm_stable(self):
        src = ROOT / "examples" / "minutes.txt"
        segs = extract(src)
        d = structure(segs, title="예시바이오 미팅", kind="minutes", client=None, chunk_chars=1600)
        md = render_md(d)
        self.assertIn("ApexBind", md)
        self.assertIn("50 billion KRW", md)
        self.assertIn("끝.", md)
        self.assertIn("2026. 7. 27", md)


if __name__ == "__main__":
    unittest.main()
