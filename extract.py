"""결정적 추출. 원본 → 세그먼트. 모델 호출 없음."""
from __future__ import annotations

import re
from dataclasses import asdict, dataclass
from pathlib import Path


DOCX = {".docx"}
PPTX = {".pptx"}
PDF = {".pdf"}
VIDEO = {".mp4", ".m4a", ".mov", ".wav", ".mp3", ".webm"}
TEXT = {".txt", ".md", ".csv"}
JSON = {".json"}


@dataclass
class Segment:
    text: str
    source: str
    location: str


def extract(path: Path) -> list[Segment]:
    path = path.expanduser().resolve()
    if not path.is_file():
        raise FileNotFoundError(f"파일 없음: {path}")
    ext = path.suffix.lower()
    if ext == ".ppt":
        raise ValueError("legacy .ppt is not supported. save as PPTX and try again.")
    if ext in PPTX:
        return extract_pptx(path)
    if ext in PDF:
        return extract_pdf(path)
    if ext in VIDEO:
        return extract_video(path)
    if ext in DOCX:
        return extract_docx(path)
    if ext in TEXT:
        return extract_text(path)
    if ext in JSON:
        return extract_text(path)
    raise ValueError(f"unsupported type: {ext}  (mp4/m4a/pptx/pdf)")


def extract_many(paths: list[Path]) -> list[Segment]:
    out: list[Segment] = []
    for p in paths:
        for seg in extract(p):
            tagged = Segment(seg.text, seg.source, f"{p.name}/{seg.location}")
            out.append(tagged)
    return out


def extract_pptx(path: Path) -> list[Segment]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    out: list[Segment] = []
    for si, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    parts.append("[표]\n" + "\n".join(rows))
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = "\n".join(p.text.strip() for p in shape.text_frame.paragraphs).strip()
            if text:
                parts.append(text)
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append("[노트]\n" + notes)
        except Exception:
            pass
        blob = "\n".join(parts).strip()
        if blob:
            out.append(Segment(blob, "ppt", f"slide-{si}"))
    return out


def extract_pdf(path: Path) -> list[Segment]:
    out: list[Segment] = []
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()
            if text:
                out.append(Segment(text, "pdf", f"page-{i}"))
    except Exception:
        out = []
    if out:
        return out
    import fitz

    doc = fitz.open(str(path))
    try:
        for i, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            if text:
                out.append(Segment(text, "pdf", f"page-{i}"))
    finally:
        doc.close()
    return out


def extract_docx(path: Path) -> list[Segment]:
    from docx import Document

    doc = Document(str(path))
    blocks: list[str] = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            blocks.append(t)
    for ti, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            rows.append(" | ".join(c.text.strip().replace("\n", " ") for c in row.cells))
        if rows:
            blocks.append(f"[표 {ti}]\n" + "\n".join(rows))
    text = "\n".join(blocks).strip()
    if not text:
        return []
    return _split_text(text, "docx")


def extract_video(path: Path) -> list[Segment]:
    sidecar = path.with_suffix(".txt")
    if sidecar.exists() and sidecar.stat().st_size > 0:
        text = sidecar.read_text(encoding="utf-8", errors="replace")
        return _transcript_segments(text, "video")
    lines = _whisper(path)
    try:
        sidecar.write_text("\n".join(lines) + "\n", encoding="utf-8")
    except OSError:
        pass
    return _transcript_segments("\n".join(lines), "video")


def _whisper(path: Path) -> list[str]:
    try:
        import mlx_whisper  # noqa: F401
        return _whisper_mlx(path)
    except ImportError:
        pass
    try:
        from faster_whisper import WhisperModel  # noqa: F401
        return _whisper_faster(path)
    except ImportError:
        pass
    raise RuntimeError(
        "no speech-to-text engine found. "
        "install `mlx-whisper av` on Mac, or `faster-whisper` on Windows/Linux."
    )


def _whisper_mlx(path: Path) -> list[str]:
    import mlx_whisper

    print(f"  transcribing with mlx-whisper ({path.name})", flush=True)
    audio = _load_audio(path)
    result = mlx_whisper.transcribe(
        audio,
        path_or_hf_repo="mlx-community/whisper-large-v3-turbo",
        language=None,
        verbose=False,
        condition_on_previous_text=False,
        no_speech_threshold=0.6,
        compression_ratio_threshold=2.4,
        logprob_threshold=-1.0,
        temperature=(0.0, 0.2, 0.4),
    )
    return [
        f"[{_fmt(seg['start'])}] {seg['text'].strip()}"
        for seg in result.get("segments") or []
        if str(seg.get("text") or "").strip()
    ]


def _whisper_faster(path: Path) -> list[str]:
    from faster_whisper import WhisperModel

    print(f"  transcribing with faster-whisper ({path.name})", flush=True)
    model = WhisperModel("small", device="auto", compute_type="int8")
    segs, _info = model.transcribe(str(path), language=None, vad_filter=True)
    out = []
    for seg in segs:
        t = (seg.text or "").strip()
        if t:
            out.append(f"[{_fmt(seg.start)}] {t}")
    return out


def _load_audio(path: Path, sr: int = 16000):
    import av
    import numpy as np

    container = av.open(str(path))
    try:
        stream = container.streams.audio[0]
        resampler = av.AudioResampler(format="s16", layout="mono", rate=sr)
        chunks = []
        for frame in container.decode(stream):
            for rf in resampler.resample(frame):
                chunks.append(rf.to_ndarray().reshape(-1))
    finally:
        container.close()
    if not chunks:
        raise RuntimeError(f"no audio track: {path.name}")
    return np.concatenate(chunks).astype(np.float32) / 32768.0


def _fmt(t) -> str:
    t = float(t or 0)
    m = int(t // 60)
    s = int(t % 60)
    return f"{m:02d}:{s:02d}"


def _transcript_segments(text: str, source: str) -> list[Segment]:
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    if not lines:
        return []
    batches: list[Segment] = []
    buf: list[str] = []
    start = "00:00"
    n = 0
    k = 1
    stamp = re.compile(r"^\[(\d{1,2}:\d{2}(?::\d{2})?)\]\s*(.*)$")
    for ln in lines:
        m = stamp.match(ln)
        body = m.group(2) if m else ln
        if m and not buf:
            start = m.group(1)
        if buf and n + len(body) > 1400:
            batches.append(Segment("\n".join(buf), source, f"{start}-{k}"))
            buf, n, k = [], 0, k + 1
            start = m.group(1) if m else start
        buf.append(ln)
        n += len(body)
    if buf:
        batches.append(Segment("\n".join(buf), source, f"{start}-{k}"))
    return batches


def extract_text(path: Path) -> list[Segment]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return _split_text(text, path.suffix.lower().lstrip("."))


def _split_text(text: str, source: str) -> list[Segment]:
    text = text.replace("\r\n", "\n").strip()
    if not text:
        return []
    chunks = re.split(r"\n(?=#{1,3}\s|\d+\.\s|[Ss]lide\s+\d+)", text)
    if len(chunks) == 1:
        paras = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
        if len(paras) > 8:
            chunks = []
            buf: list[str] = []
            n = 0
            for p in paras:
                buf.append(p)
                n += len(p)
                if n >= 1400:
                    chunks.append("\n\n".join(buf))
                    buf, n = [], 0
            if buf:
                chunks.append("\n\n".join(buf))
        else:
            chunks = [text]
    out = []
    for i, c in enumerate(chunks, 1):
        c = c.strip()
        if c:
            out.append(Segment(c, source, f"part-{i}"))
    return out


def hangul_ratio(text: str) -> float:
    letters = [c for c in text if c.isalpha() or "\uac00" <= c <= "\ud7a3"]
    if not letters:
        return 0.0
    han = sum(1 for c in letters if "\uac00" <= c <= "\ud7a3")
    return han / len(letters)


def collect_numbers(text: str) -> list[str]:
    pat = re.compile(
        r"(?:\d{1,3}(?:,\d{3})+|\d+(?:\.\d+)?)(?:\s*(?:억원|만원|천원|달러|엔|원|억|만|천|%|명|건|개|년|월|일|시간|분))?"
    )
    seen: list[str] = []
    for m in pat.findall(text):
        s = m.strip()
        if s not in seen:
            seen.append(s)
    return seen


def as_dicts(segments: list[Segment]) -> list[dict]:
    return [asdict(s) for s in segments]
